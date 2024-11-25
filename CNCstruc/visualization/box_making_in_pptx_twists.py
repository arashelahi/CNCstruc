from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import numpy as np
import traj_reader as trj

import CNC_class as cnc
import cnc_analysis_utils
from pptx.util import Inches, Pt

from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Function to add a rectangle to a slide
def add_rectangle(slide, left, top, width, height,rgb_color,domain):


    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,  # Shape type: Rectangle
        left, top, width, height  # Dimensions and position
    )
    # Set solid fill color and transparency
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255,255,255)  # Example: white fill
    # shape.fill.fore_color.rgb = rgb_color  # Example: white fill
    shape.fill.transparency = 91        # 91% transparency

    # Set line color and width
    shape.line.color.rgb = RGBColor(0,0,0)             # Black line border
    shape.line.width = Pt(1)  # Set line width to 1 point
    if domain == 'exterior':
        shape.line.fill.transparency = 80 
    if shape.shadow:
        shape.shadow.inherit = False
        shape.shadow.visible = False
    return shape

def create_presentation(CNC_group):
    height = Inches(0.5)
    width = Inches(1.2)
    ppt = Presentation()  # Load an existing presentation

    # rect = add_rectangle(slide, Inches(4.98), Inches(2.19), height, width)
    slide_layout = ppt.slide_layouts[1]
    slide1 = ppt.slides.add_slide(slide_layout)
    slide2 = ppt.slides.add_slide(slide_layout)

    ## make the    # Lists to hold shapes by domain
    exterior_shapes = []
    interior_shapes = []
    # Step 1: make the powerpoint
    # If you have a specific template with boxes, load it, else create a new one

    # Step 2: The directory of the images
    Interior_Image_dir = './twist_data/Images/interior/'
    Exterior_Image_dir = './twist_data/Images/exterior/'

    ppt_dir = './twist_data/Images/'
    # Step 3: Embed each plot into each box in the slide

    horizontal_spacing = Inches(1.27)  # Space between images in the same layer
    vertical_spacing = Inches(0.56)    # Space between layers
    initial_offset_x = Inches(4) # Initial horizontal offset
    initial_offset_y = Inches(1) # Initial vertical offset
    plot_to_box_offset_x = Inches(0.19)
    plot_to_box_offset_y = Inches(0.03)
    top_lays_x_diff = Inches(-0.69)
    bot_lays_x_diff = Inches(0.58)


    for layer_index, layer in enumerate(CNC_group.layer_vec):
        initial_offset_x = initial_offset_x if layer_index == 0 \
                    else initial_offset_x + top_lays_x_diff if layer_index < 6 \
                    else initial_offset_x + bot_lays_x_diff if layer_index >= 6 else initial_offset_x

        chain_number_vec = CNC_group.layers[layer]
        top = initial_offset_y + layer_index * vertical_spacing
        for chain_iter,chain_number in enumerate(chain_number_vec):
            left = initial_offset_x + chain_iter * horizontal_spacing
            if len(chain_number_vec)<=2 or chain_iter == 0 or chain_iter == len(chain_number_vec)-1:
                domain = 'exterior'
                rgb_color = RGBColor(0, 0, 0)
                file_ext = '_ends'
                plot_file_name = Exterior_Image_dir+'%s_ch%s_%s_twist_boxplot.png' % (layer, chain_iter if chain_iter==0 else 1 , domain)
            else:
                domain = 'interior'
                rgb_color = RGBColor(192, 0, 0)
                file_ext = ''
                plot_file_name = Interior_Image_dir+'%s_ch%s_%s_twist_boxplot.png' % (layer, chain_iter-1, domain)


            add_rectangle(slide1, left, top, width, height,rgb_color,domain)
            slide2.shapes.add_picture(plot_file_name, left + plot_to_box_offset_x, top + plot_to_box_offset_y)


    ppt.save(ppt_dir+'boxes.pptx')

file_dir = './simulation_traj_topol/'
gro_file = file_dir+'solute.gro'

init_domain = 'exterior' ## 'interior' or 'exterior
Data = trj.gro_reader(gro_file)
CNC_group = cnc.CNC_analys(Data,init_domain)
create_presentation(CNC_group)






# for i, (pos, size) in enumerate(zip(positions, sizes)):
#     left, top = Inches(pos[0]), Inches(pos[1])  # Position of each image
#     # width, height = Inches(size[0]), Inches(size[1])  # Size of each image
#     slide.shapes.add_picture(f'plot_{i+1}.png', left, top)
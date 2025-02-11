from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import numpy as np
import traj_reader as trj

import CNC_class as cnc
import cnc_analysis_utils

file_dir = './simulation_traj_topol/'
gro_file = file_dir+'solute.gro'

domain = 'interior' ## 'interior' or 'exterior
file_ext = '' ## if the interior chains are needed
if domain!='interior':
    file_ext = '_ends' ## if the exterior chains are needed
Data = trj.gro_reader(gro_file)
CNC_group = cnc.CNC_analys(Data,domain)

# Step 1: Load your presentation or create a new one
# If you have a specific template with boxes, load it, else create a new one
ppt = Presentation('./tg_data/Images/test_pptx.pptx')  # Load an existing presentation

# Step 2: The directory of the images
Image_dir = './tg_data/Images/interior/'
# Step 3: Embed each plot into each box in the slide

horizontal_spacing = Inches(1.27)  # Space between images in the same layer
vertical_spacing = Inches(0.56)    # Space between layers
# initial_offset = Inches(1)      # Starting position for the first image
initial_offset_x = Inches(4.96) # Initial horizontal offset
initial_offset_y = Inches(2.19) # Initial vertical offset
slide = ppt.slides[2]  # Assuming you want to insert into the first slide

for layer_index, layer in enumerate(CNC_group.layer_vec):
    initial_offset_x = initial_offset_x if layer_index == 0 \
                  else initial_offset_x + Inches(-0.62) if layer_index < 4 \
                  else initial_offset_x + Inches(0.50) if layer_index >= 4 else initial_offset_x
    if domain=='interior': 
        chain_number_vec = CNC_group.layers[layer][1:-1] # For interior
    else:
        chain_number_vec = [CNC_group.layers[layer][0], CNC_group.layers[layer][-1]] if len(CNC_group.layers[layer]) > 1\
            else [CNC_group.layers[layer][0]] # for the exterior chains
    top = initial_offset_y + layer_index * vertical_spacing
    for chain_iter,chain_number in enumerate(chain_number_vec):
        left = initial_offset_x + chain_iter * horizontal_spacing
        # left, top = Inches(1), Inches(1)  # Position of each image
        slide.shapes.add_picture(Image_dir+f'{layer}_ch{chain_iter}_tgs{file_ext}.png', left, top)


# Step 4: Save the updated presentation
ppt.save(Image_dir+'updated_presentation.pptx')


# for i, (pos, size) in enumerate(zip(positions, sizes)):
#     left, top = Inches(pos[0]), Inches(pos[1])  # Position of each image
#     # width, height = Inches(size[0]), Inches(size[1])  # Size of each image
#     slide.shapes.add_picture(f'plot_{i+1}.png', left, top)